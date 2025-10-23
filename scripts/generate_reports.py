from __future__ import annotations

import csv
from pathlib import Path
import sqlite3
import pandas as pd

# Optional libraries for PDF and PPTX
try:
    from reportlab.lib.pagesizes import LETTER
    from reportlab.pdfgen import canvas
    HAVE_PDF = True
except Exception:
    HAVE_PDF = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    HAVE_PPTX = True
except Exception:
    HAVE_PPTX = False

try:
    from scripts.config import ROOT, FILES, REPORTS_DIR, FIGURES_DIR
except ModuleNotFoundError:
    from config import ROOT, FILES, REPORTS_DIR, FIGURES_DIR

DB_PATH = ROOT / "data" / "ai_talks.sqlite"


def generate_data_dictionary(out_path: Path) -> None:
    rows = []
    for name, path in FILES.items():
        p = Path(path)
        if not p.exists():
            continue
        df = pd.read_csv(p, nrows=100)
        for col in df.columns:
            dtype = str(df[col].dtype)
            rows.append({
                "table": name,
                "column": col,
                "dtype": dtype,
                "description": "",
            })
    out_path.parent.mkdir(parents=True, exist_ok=True)
    pd.DataFrame(rows).to_csv(out_path, index=False)


def generate_executive_summary_pdf(out_path: Path) -> None:
    if not HAVE_PDF:
        print("[WARN] reportlab not installed; skipping PDF.")
        return
    kpis_csv = REPORTS_DIR / "kpis.csv"
    kpis = pd.read_csv(kpis_csv) if kpis_csv.exists() else pd.DataFrame()

    c = canvas.Canvas(str(out_path), pagesize=LETTER)
    width, height = LETTER
    y = height - 72
    c.setFont("Helvetica-Bold", 16)
    c.drawString(72, y, "AI Talks Campaign Executive Summary")
    y -= 24
    c.setFont("Helvetica", 10)
    if not kpis.empty:
        for key in kpis.columns:
            y -= 14
            c.drawString(72, y, f"{key}: {kpis.iloc[0][key]}")
    else:
        y -= 14
        c.drawString(72, y, "KPIs not available. Run eda_youtube.py to generate kpis.csv.")

    y -= 28
    c.setFont("Helvetica-Bold", 12)
    c.drawString(72, y, "Key Visuals")
    y -= 12
    # Place a few core figures if present
    figures = [
        FIGURES_DIR / "top_videos_by_views.png",
        FIGURES_DIR / "traffic_sources.png",
        FIGURES_DIR / "top_countries.png",
        FIGURES_DIR / "subs_over_time.png",
        FIGURES_DIR / "subscriber_breakdown.png",
    ]
    for fig in figures:
        if fig.exists():
            y -= 200
            try:
                c.drawImage(str(fig), 72, max(y, 72), width=width-144, height=180, preserveAspectRatio=True, anchor='n')
            except Exception:
                pass
            if y < 100:
                c.showPage()
                y = height - 72
    c.showPage()
    c.save()


def generate_insights_pptx(out_path: Path) -> None:
    if not HAVE_PPTX:
        print("[WARN] python-pptx not installed; skipping PPTX.")
        return
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "AI Talks Campaign Insights"
    slide.placeholders[1].text = "Auto-generated deck with KPIs and key visuals"

    # KPIs slide
    kpis_csv = REPORTS_DIR / "kpis.csv"
    if kpis_csv.exists():
        df = pd.read_csv(kpis_csv)
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Key KPIs"
        body = slide.placeholders[1].text_frame
        body.text = ""
        for col in df.columns:
            body.add_paragraph().text = f"{col}: {df.iloc[0][col]}"

    # Visuals slides
    for title, filename in [
        ("Top Videos by Views", "top_videos_by_views.png"),
        ("Traffic Sources", "traffic_sources.png"),
        ("Top Countries", "top_countries.png"),
        ("Subscribers Over Time", "subs_over_time.png"),
        ("Subscriber Breakdown", "subscriber_breakdown.png"),
    ]:
        path = FIGURES_DIR / filename
        if not path.exists():
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        slide.shapes.add_picture(str(path), left, top, width=width)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main() -> None:
    REPORTS_DIR.mkdir(parents=True, exist_ok=True)
    # Data dictionary
    generate_data_dictionary(REPORTS_DIR / "data_dictionary.csv")
    # PDF summary
    generate_executive_summary_pdf(REPORTS_DIR / "executive_summary.pdf")
    # PPTX deck
    generate_insights_pptx(REPORTS_DIR / "insights_presentation.pptx")
    print("Reports generated under /reports")


if __name__ == "__main__":
    main()
